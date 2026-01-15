import os
import sqlite3
from datetime import datetime
from flask import Flask, render_template, request, redirect, url_for, jsonify, send_file, flash

try:
    from pptx import Presentation
    from pptx.util import Pt
    PPTX_OK = True
except Exception:
    PPTX_OK = False

APP_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(APP_DIR, "gpse.db")
PPT_DIR = os.path.join(APP_DIR, "output_ppt")
os.makedirs(PPT_DIR, exist_ok=True)

app = Flask(__name__)
app.secret_key = "gpse-v2-plus"

def conn():
    c = sqlite3.connect(DB_PATH)
    c.row_factory = sqlite3.Row
    return c

def now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")

def init_db():
    with conn() as c:
        c.execute("""CREATE TABLE IF NOT EXISTS reports (
            report_id TEXT PRIMARY KEY,
            project TEXT NOT NULL,
            week TEXT NOT NULL,
            owner TEXT NOT NULL,
            report_type TEXT NOT NULL,
            status TEXT NOT NULL,
            storage_url TEXT,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )""")
        c.execute("""CREATE TABLE IF NOT EXISTS versions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            report_id TEXT NOT NULL,
            version_no INTEGER NOT NULL,
            notes TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(report_id) REFERENCES reports(report_id)
        )""")
        c.execute("""CREATE TABLE IF NOT EXISTS kpis (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            report_id TEXT NOT NULL,
            sla REAL,
            p1_incidents INTEGER,
            mttr_minutes INTEGER,
            risk_count INTEGER,
            rag TEXT,
            created_at TEXT NOT NULL,
            FOREIGN KEY(report_id) REFERENCES reports(report_id)
        )""")

        c.execute("""CREATE TABLE IF NOT EXISTS departments (
            dept_id INTEGER PRIMARY KEY AUTOINCREMENT,
            dept_name TEXT UNIQUE NOT NULL
        )""")

        c.execute("""CREATE TABLE IF NOT EXISTS kpi_master (
            kpi_id INTEGER PRIMARY KEY AUTOINCREMENT,
            dept_id INTEGER NOT NULL,
            section TEXT NOT NULL,
            kpi_key TEXT NOT NULL,
            kpi_name TEXT NOT NULL,
            formula_display TEXT NOT NULL,
            description TEXT,
            calculation_notes TEXT,
            green_rule TEXT,
            amber_rule TEXT,
            red_rule TEXT,
            owner_team TEXT,
            updated_at TEXT NOT NULL,
            UNIQUE(dept_id, kpi_key),
            FOREIGN KEY(dept_id) REFERENCES departments(dept_id)
        )""")

        if c.execute("SELECT COUNT(*) AS n FROM departments").fetchone()["n"] == 0:
            seed_departments(c)
        if c.execute("SELECT COUNT(*) AS n FROM kpi_master").fetchone()["n"] == 0:
            seed_kpi_library(c)
        if c.execute("SELECT COUNT(*) AS n FROM reports").fetchone()["n"] == 0:
            seed_reports_demo(c)

def seed_departments(c):
    for d in ["NFPE", "INC (Incident)", "CRI", "Service Desk", "GPSE Ops"]:
        c.execute("INSERT INTO departments(dept_name) VALUES(?)", (d,))

def _dept_id(c, dept_name: str) -> int:
    return c.execute("SELECT dept_id FROM departments WHERE dept_name=?", (dept_name,)).fetchone()["dept_id"]

def seed_kpi_library(c):
    def upsert(dept_name, section, kpi_key, kpi_name, formula, desc="", notes="", g="", a="", r="", owner="GPSE"):
        did = _dept_id(c, dept_name)
        c.execute("""
        INSERT INTO kpi_master(dept_id, section, kpi_key, kpi_name, formula_display, description, calculation_notes,
                               green_rule, amber_rule, red_rule, owner_team, updated_at)
        VALUES(?,?,?,?,?,?,?,?,?,?,?,?)
        ON CONFLICT(dept_id, kpi_key) DO UPDATE SET
          section=excluded.section,
          kpi_name=excluded.kpi_name,
          formula_display=excluded.formula_display,
          description=excluded.description,
          calculation_notes=excluded.calculation_notes,
          green_rule=excluded.green_rule,
          amber_rule=excluded.amber_rule,
          red_rule=excluded.red_rule,
          owner_team=excluded.owner_team,
          updated_at=excluded.updated_at
        """, (did, section, kpi_key, kpi_name, formula, desc, notes, g, a, r, owner, now()))

    upsert("NFPE", "NFPE", "NFPE_CRI_SEV2_RATE", "CRI SEV2 Rate",
           "CRI SEV2 Rate = (Count of CRI tickets tagged SEV2) / (Total CRI tickets) × 100",
           desc="Measures concentration of SEV2 within CRI tickets.",
           notes="Filter: selected reporting period; CRI scope defined by NFPE taxonomy.",
           g="Green: < 2%", a="Amber: 2%–5%", r="Red: > 5%")

    upsert("NFPE", "NFPE", "NFPE_EXCEPTION_RATE", "Exception Rate",
           "Exception Rate = (Exception count) / (Total cases processed) × 100",
           desc="Shows operational exceptions relative to volume.",
           notes="Define exceptions per NFPE policy; exclude training/test cases.",
           g="Green: < 0.5%", a="Amber: 0.5%–1.0%", r="Red: > 1.0%")

    upsert("INC (Incident)", "Incident Mgmt", "INC_P1_COUNT", "P1 Incidents",
           "P1 = count(priority='P1') in reporting period",
           desc="Number of Priority-1 incidents created in the period.",
           notes="Use incident creation timestamp; exclude duplicates/cancelled.",
           g="Green: 0", a="Amber: 1–2", r="Red: ≥ 3")

    upsert("INC (Incident)", "Incident Mgmt", "INC_MTTR_P1", "P1 MTTR (minutes)",
           "MTTR(P1) = Avg(Resolved Time − Opened Time) for Priority=P1",
           desc="Mean time to resolve P1 incidents.",
           notes="Use resolved incidents only; exclude vendor-hold or paused time if policy requires.",
           g="Green: ≤ 45", a="Amber: 46–90", r="Red: > 90")

    upsert("CRI", "CRI", "CRI_BACKLOG", "CRI Backlog",
           "Backlog = count(CRI items where status in {Open, In Progress})",
           desc="Open backlog items under CRI scope.",
           notes="Backlog definition agreed with CRI ops; snapshot at end of period.",
           g="Green: ≤ 10", a="Amber: 11–25", r="Red: > 25")

    upsert("Service Desk", "Service Desk", "SD_SLA", "SLA Compliance",
           "SLA% = (Success / Total) × 100",
           desc="Share of service desk requests meeting SLA.",
           notes="Define success as met SLA within policy; period based on ticket closed date.",
           g="Green: ≥ 99.0%", a="Amber: 97.0%–98.99%", r="Red: < 97.0%")

    upsert("GPSE Ops", "Ops", "OPS_ACTIVE_RISKS", "Active Risks",
           "Active Risks = count(risks where status='Active')",
           desc="Count of open/active risks tracked by GPSE Ops.",
           notes="Risk register maintained weekly; treat overdue mitigations as active.",
           g="Green: ≤ 2", a="Amber: 3–5", r="Red: > 5")

def seed_reports_demo(c):
    demo_reports = [
        ("R001", "Alpha", "W35", "GPSE1", "Weekly", "Final", "https://sharepoint.example/Alpha_W35"),
        ("R002", "Beta",  "W35", "GPSE2", "Incident", "Draft", "https://sharepoint.example/Beta_W35"),
        ("R003", "Alpha", "W36", "GPSE1", "Weekly", "Draft", "https://sharepoint.example/Alpha_W36"),
        ("R004", "Gamma", "W35", "GPSE3", "Weekly", "Final", "https://sharepoint.example/Gamma_W35"),
    ]
    for r in demo_reports:
        c.execute("""INSERT INTO reports
        (report_id, project, week, owner, report_type, status, storage_url, created_at, updated_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""", (*r, now(), now()))

    demo_versions = [
        ("R001", 1, "Initial weekly report draft."),
        ("R001", 2, "Added KPI clarifications and updated risks."),
        ("R002", 1, "Incident summary created."),
        ("R003", 1, "New week created; waiting for final KPIs."),
    ]
    for report_id, vno, notes in demo_versions:
        c.execute("INSERT INTO versions(report_id, version_no, notes, created_at) VALUES(?,?,?,?)",
                  (report_id, vno, notes, now()))

    demo_kpis = [
        ("R001", 99.2, 1, 45, 2, "Green"),
        ("R002", 97.1, 3, 80, 5, "Amber"),
        ("R003", 98.7, 2, 60, 3, "Amber"),
        ("R004", 99.6, 0, 30, 1, "Green"),
    ]
    for report_id, sla, p1, mttr, risks, rag in demo_kpis:
        c.execute("""INSERT INTO kpis(report_id, sla, p1_incidents, mttr_minutes, risk_count, rag, created_at)
                     VALUES(?,?,?,?,?,?,?)""", (report_id, sla, p1, mttr, risks, rag, now()))

def summary_cards(c):
    return c.execute("""SELECT
          AVG(sla) AS avg_sla,
          SUM(p1_incidents) AS total_p1,
          AVG(mttr_minutes) AS avg_mttr,
          SUM(risk_count) AS total_risks
        FROM kpis""").fetchone()

@app.route("/")
def root():
    return redirect(url_for("dashboard"))

@app.route("/dashboard")
def dashboard():
    init_db()
    q = request.args.get("q", "").strip().lower()
    project = request.args.get("project", "").strip()
    week = request.args.get("week", "").strip()
    owner = request.args.get("owner", "").strip()
    report_type = request.args.get("report_type", "").strip()

    sql = "SELECT * FROM reports WHERE 1=1"
    params = []
    if q:
        sql += " AND (lower(report_id) LIKE ? OR lower(project) LIKE ? OR lower(owner) LIKE ? OR lower(report_type) LIKE ?)"
        params += [f"%{q}%", f"%{q}%", f"%{q}%", f"%{q}%"]
    if project:
        sql += " AND project=?"; params.append(project)
    if week:
        sql += " AND week=?"; params.append(week)
    if owner:
        sql += " AND owner=?"; params.append(owner)
    if report_type:
        sql += " AND report_type=?"; params.append(report_type)
    sql += " ORDER BY updated_at DESC"

    with conn() as c:
        reports = c.execute(sql, params).fetchall()
        projects = [r["project"] for r in c.execute("SELECT DISTINCT project FROM reports ORDER BY project").fetchall()]
        weeks = [r["week"] for r in c.execute("SELECT DISTINCT week FROM reports ORDER BY week DESC").fetchall()]
        owners = [r["owner"] for r in c.execute("SELECT DISTINCT owner FROM reports ORDER BY owner").fetchall()]
        types = [r["report_type"] for r in c.execute("SELECT DISTINCT report_type FROM reports ORDER BY report_type").fetchall()]
        cards = summary_cards(c)

    return render_template("dashboard.html", active="dashboard",
                           reports=reports, projects=projects, weeks=weeks, owners=owners, types=types, cards=cards)

@app.route("/create", methods=["GET","POST"])
def create():
    init_db()
    if request.method == "POST":
        report_id = request.form.get("report_id","").strip()
        project = request.form.get("project","").strip()
        week = request.form.get("week","").strip()
        owner = request.form.get("owner","").strip()
        report_type = request.form.get("report_type","").strip()
        status = request.form.get("status","Draft").strip()
        storage_url = request.form.get("storage_url","").strip()

        if not all([report_id, project, week, owner, report_type, status]):
            flash("All fields except Storage URL are required.", "danger")
            return redirect(url_for("create"))

        with conn() as c:
            exists = c.execute("SELECT 1 FROM reports WHERE report_id=?", (report_id,)).fetchone()
            if exists:
                flash("Report ID already exists. Use a new ID.", "danger")
                return redirect(url_for("create"))
            c.execute("""INSERT INTO reports
                (report_id, project, week, owner, report_type, status, storage_url, created_at, updated_at)
                VALUES (?,?,?,?,?,?,?,?,?)""", (report_id, project, week, owner, report_type, status, storage_url, now(), now()))
        flash("Report created.", "success")
        return redirect(url_for("report_detail", report_id=report_id))
    return render_template("create.html", active="dashboard")

@app.route("/report/<report_id>")
def report_detail(report_id):
    init_db()
    with conn() as c:
        report = c.execute("SELECT * FROM reports WHERE report_id=?", (report_id,)).fetchone()
        if not report:
            return "Report not found", 404
        versions = c.execute("SELECT * FROM versions WHERE report_id=? ORDER BY version_no DESC", (report_id,)).fetchall()
        kpi_latest = c.execute("SELECT * FROM kpis WHERE report_id=? ORDER BY created_at DESC LIMIT 1", (report_id,)).fetchone()
    return render_template("report_detail.html", active="dashboard", report=report, versions=versions, kpi=kpi_latest)

@app.route("/add_version/<report_id>", methods=["POST"])
def add_version(report_id):
    init_db()
    notes = request.form.get("notes","").strip()
    if not notes:
        flash("Version notes are required.", "danger")
        return redirect(url_for("report_detail", report_id=report_id))
    with conn() as c:
        mx = c.execute("SELECT COALESCE(MAX(version_no),0) AS mx FROM versions WHERE report_id=?", (report_id,)).fetchone()["mx"]
        vno = int(mx) + 1
        c.execute("INSERT INTO versions(report_id, version_no, notes, created_at) VALUES(?,?,?,?)",
                  (report_id, vno, notes, now()))
        c.execute("UPDATE reports SET updated_at=? WHERE report_id=?", (now(), report_id))
    flash(f"Saved version v{vno}.", "success")
    return redirect(url_for("report_detail", report_id=report_id))

@app.route("/add_kpi/<report_id>", methods=["POST"])
def add_kpi(report_id):
    init_db()
    def _to_int(v, default=0):
        try: return int(v)
        except Exception: return default
    def _to_float(v, default=0.0):
        try: return float(v)
        except Exception: return default

    sla = _to_float(request.form.get("sla","0"))
    p1 = _to_int(request.form.get("p1_incidents","0"))
    mttr = _to_int(request.form.get("mttr_minutes","0"))
    risks = _to_int(request.form.get("risk_count","0"))
    rag = request.form.get("rag","Green").strip()

    with conn() as c:
        c.execute("""INSERT INTO kpis(report_id, sla, p1_incidents, mttr_minutes, risk_count, rag, created_at)
                     VALUES(?,?,?,?,?,?,?)""", (report_id, sla, p1, mttr, risks, rag, now()))
        c.execute("UPDATE reports SET updated_at=? WHERE report_id=?", (now(), report_id))
    flash("KPI snapshot saved.", "success")
    return redirect(url_for("report_detail", report_id=report_id))

def _ppt_add_bullets(slide, title, bullets):
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

@app.route("/generate_ppt/<report_id>", methods=["POST"])
def generate_ppt(report_id):
    init_db()
    if not PPTX_OK:
        flash("python-pptx not installed. Ask IT to allow install or remove PPT feature for demo.", "warning")
        return redirect(url_for("report_detail", report_id=report_id))

    with conn() as c:
        report = c.execute("SELECT * FROM reports WHERE report_id=?", (report_id,)).fetchone()
        if not report:
            flash("Report not found.", "danger")
            return redirect(url_for("dashboard"))
        kpi = c.execute("SELECT * FROM kpis WHERE report_id=? ORDER BY created_at DESC LIMIT 1", (report_id,)).fetchone()
        versions = c.execute("SELECT * FROM versions WHERE report_id=? ORDER BY version_no DESC LIMIT 5", (report_id,)).fetchall()

    prs = Presentation()
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = f"{report['project']} — {report['week']} ({report['report_type']})"
    s0.placeholders[1].text = f"Report ID: {report['report_id']} | Owner: {report['owner']} | Status: {report['status']}\nGenerated: {now()}"

    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    bullets = [f"SLA: {kpi['sla']}%", f"P1 Incidents: {kpi['p1_incidents']}", f"MTTR: {kpi['mttr_minutes']} min",
               f"Risk Count: {kpi['risk_count']}", f"RAG Status: {kpi['rag']}"] if kpi else ["No KPI snapshot found yet."]
    _ppt_add_bullets(s1, "KPIs (Latest Snapshot)", bullets)

    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    bullets = [f"v{v['version_no']}: {v['notes']}" for v in versions] if versions else ["No versions recorded yet."]
    _ppt_add_bullets(s2, "Recent Versions / Notes", bullets)

    filename = f"{report_id}_{report['project']}_{report['week']}.pptx".replace(" ", "_")
    out_path = os.path.join(PPT_DIR, filename)
    prs.save(out_path)

    flash("PPT generated successfully. Check output_ppt/ folder.", "success")
    return redirect(url_for("report_detail", report_id=report_id))

@app.route("/kpi-library")
def kpi_library():
    init_db()
    return render_template("kpi_library.html", active="kpi_library")

@app.route("/api/departments")
def api_departments():
    init_db()
    with conn() as c:
        rows = c.execute("SELECT dept_id, dept_name FROM departments ORDER BY dept_name").fetchall()
    return jsonify([dict(r) for r in rows])

@app.route("/api/kpis_list")
def api_kpis_list():
    init_db()
    dept_id = request.args.get("dept_id","").strip()
    section = request.args.get("section","").strip()
    search = request.args.get("search","").strip().lower()

    sql = "SELECT kpi_id, section, kpi_key, kpi_name, updated_at FROM kpi_master WHERE 1=1"
    params = []
    if dept_id:
        sql += " AND dept_id=?"; params.append(dept_id)
    if section:
        sql += " AND section=?"; params.append(section)
    if search:
        sql += " AND (lower(kpi_key) LIKE ? OR lower(kpi_name) LIKE ? OR lower(section) LIKE ?)"
        params += [f"%{search}%", f"%{search}%", f"%{search}%"]
    sql += " ORDER BY section ASC, kpi_name ASC"

    with conn() as c:
        rows = c.execute(sql, params).fetchall()
    return jsonify([dict(r) for r in rows])

@app.route("/api/kpi/<int:kpi_id>")
def api_kpi_detail(kpi_id):
    init_db()
    with conn() as c:
        row = c.execute("""
        SELECT km.*, d.dept_name
        FROM kpi_master km
        JOIN departments d ON d.dept_id = km.dept_id
        WHERE km.kpi_id=?
        """, (kpi_id,)).fetchone()
    if not row:
        return jsonify({"error":"KPI not found"}), 404
    return jsonify(dict(row))

@app.route("/api/kpi_master")
def api_kpi_master():
    init_db()
    with conn() as c:
        rows = c.execute("""
        SELECT d.dept_name, km.section, km.kpi_key, km.kpi_name, km.formula_display,
               km.description, km.calculation_notes, km.green_rule, km.amber_rule, km.red_rule,
               km.owner_team, km.updated_at
        FROM kpi_master km
        JOIN departments d ON d.dept_id = km.dept_id
        ORDER BY d.dept_name, km.section, km.kpi_name
        """).fetchall()
    return jsonify([dict(r) for r in rows])

@app.route("/utilities")
def utilities():
    init_db()
    return render_template("utilities.html", active="utilities")

@app.route("/assistant")
def assistant():
    init_db()
    q = request.args.get("q","").strip()
    results = []
    if q:
        ql = q.lower()
        with conn() as c:
            results = c.execute("""
            SELECT km.kpi_name, km.kpi_key, km.section, km.formula_display, d.dept_name
            FROM kpi_master km
            JOIN departments d ON d.dept_id = km.dept_id
            WHERE lower(km.kpi_name) LIKE ? OR lower(km.kpi_key) LIKE ? OR lower(km.section) LIKE ? OR lower(d.dept_name) LIKE ?
            ORDER BY d.dept_name, km.section, km.kpi_name
            LIMIT 50
            """, (f"%{ql}%", f"%{ql}%", f"%{ql}%", f"%{ql}%")).fetchall()
    return render_template("assistant.html", active="assistant", q=q, results=results)

@app.route("/api/reports")
def api_reports():
    init_db()
    with conn() as c:
        rows = c.execute("SELECT * FROM reports").fetchall()
    out = []
    for r in rows:
        d = dict(r)
        d["hub_url"] = f"http://127.0.0.1:5000/report/{d['report_id']}"
        out.append(d)
    return jsonify(out)

@app.route("/api/kpis")
def api_kpis():
    init_db()
    with conn() as c:
        rows = c.execute("""SELECT
              r.report_id, r.project, r.week, r.owner, r.report_type, r.status,
              k.sla, k.p1_incidents, k.mttr_minutes, k.risk_count, k.rag, k.created_at
            FROM reports r
            LEFT JOIN (
              SELECT kk.*
              FROM kpis kk
              INNER JOIN (SELECT report_id, MAX(created_at) AS mc FROM kpis GROUP BY report_id) latest
              ON latest.report_id = kk.report_id AND latest.mc = kk.created_at
            ) k
            ON k.report_id = r.report_id""").fetchall()
    out = []
    for r in rows:
        d = dict(r)
        d["hub_url"] = f"http://127.0.0.1:5000/report/{d['report_id']}"
        out.append(d)
    return jsonify(out)

@app.route("/export/kpis.csv")
def export_kpis_csv():
    init_db()
    with conn() as c:
        rows = c.execute("""SELECT
              r.report_id, r.project, r.week, r.owner, r.report_type, r.status,
              k.sla, k.p1_incidents, k.mttr_minutes, k.risk_count, k.rag, k.created_at
            FROM reports r
            LEFT JOIN (
              SELECT kk.*
              FROM kpis kk
              INNER JOIN (SELECT report_id, MAX(created_at) AS mc FROM kpis GROUP BY report_id) latest
              ON latest.report_id = kk.report_id AND latest.mc = kk.created_at
            ) k
            ON k.report_id = r.report_id
            ORDER BY r.week DESC, r.project ASC""").fetchall()

    import io, csv
    output = io.StringIO()
    w = csv.writer(output)
    w.writerow(["report_id","project","week","owner","report_type","status","sla","p1_incidents","mttr_minutes","risk_count","rag","kpi_updated_at","hub_url"])
    for r in rows:
        d = dict(r)
        w.writerow([d.get("report_id",""), d.get("project",""), d.get("week",""), d.get("owner",""),
                    d.get("report_type",""), d.get("status",""), d.get("sla",""), d.get("p1_incidents",""),
                    d.get("mttr_minutes",""), d.get("risk_count",""), d.get("rag",""), d.get("created_at",""),
                    f"http://127.0.0.1:5000/report/{d.get('report_id','')}"])
    mem = io.BytesIO(output.getvalue().encode("utf-8"))
    return send_file(mem, mimetype="text/csv", as_attachment=True, download_name="kpis_export.csv")

@app.route("/export/kpi_library.csv")
def export_kpi_master_csv():
    init_db()
    with conn() as c:
        rows = c.execute("""
        SELECT d.dept_name, km.section, km.kpi_key, km.kpi_name, km.formula_display,
               km.description, km.calculation_notes, km.green_rule, km.amber_rule, km.red_rule,
               km.owner_team, km.updated_at
        FROM kpi_master km
        JOIN departments d ON d.dept_id = km.dept_id
        ORDER BY d.dept_name, km.section, km.kpi_name
        """).fetchall()

    import io, csv
    output = io.StringIO()
    w = csv.writer(output)
    w.writerow(["dept_name","section","kpi_key","kpi_name","formula_display","description","calculation_notes","green_rule","amber_rule","red_rule","owner_team","updated_at"])
    for r in rows:
        d = dict(r)
        w.writerow([d.get("dept_name",""), d.get("section",""), d.get("kpi_key",""), d.get("kpi_name",""),
                    d.get("formula_display",""), d.get("description",""), d.get("calculation_notes",""),
                    d.get("green_rule",""), d.get("amber_rule",""), d.get("red_rule",""),
                    d.get("owner_team",""), d.get("updated_at","")])
    mem = io.BytesIO(output.getvalue().encode("utf-8"))
    return send_file(mem, mimetype="text/csv", as_attachment=True, download_name="kpi_library_export.csv")

@app.route("/admin/reset-demo", methods=["POST"])
def reset_demo():
    if os.path.exists(DB_PATH):
        os.remove(DB_PATH)
    init_db()
    flash("Demo database reset and re-seeded.", "success")
    return redirect(url_for("dashboard"))

if __name__ == "__main__":
    init_db()
    app.run(host="127.0.0.1", port=5000, debug=True, use_reloader=False)
